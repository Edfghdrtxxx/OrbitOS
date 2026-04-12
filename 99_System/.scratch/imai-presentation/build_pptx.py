"""
Build 5-slide presentation for meeting with Prof. Nobuaki Imai (UTokyo CNS).
Presenter: ZhiHeng Hu, IMP (CAS) Master's student.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Constants ──────────────────────────────────────────────────────────────
SLIDE_WIDTH = Inches(13.333)  # 16:9
SLIDE_HEIGHT = Inches(7.5)
FONT_TITLE = "Calibri"
FONT_BODY = "Calibri"
COLOR_TITLE = RGBColor(0x1B, 0x3A, 0x5C)       # dark navy
COLOR_SUBTITLE = RGBColor(0x4A, 0x4A, 0x4A)     # medium gray
COLOR_BODY = RGBColor(0x2D, 0x2D, 0x2D)         # near-black
COLOR_ACCENT = RGBColor(0x00, 0x6B, 0xA6)       # blue accent
COLOR_HIGHLIGHT = RGBColor(0xC0, 0x39, 0x2B)     # red for key numbers
COLOR_BG_STRIPE = RGBColor(0x1B, 0x3A, 0x5C)    # navy header stripe
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

OUT_DIR = r"D:\obsidian\OrbitOS\99_System\.scratch\imai-presentation"
OUT_FILE = os.path.join(OUT_DIR, "Imai_Meeting_20260413.pptx")

# ── Figure paths (will try to embed PNGs) ─────────────────────────────────
FIGURES = {
    "mate_schema": r"D:\Something\research\MATE-Automation-V4\10_Papers-Thesis\Physics_Informed\figures\mate_geometry_schema.png",
    "example_tracks": r"D:\Something\research\MATE-Automation-V4\10_Papers-Thesis\Physics_Informed\figures\example_tracks.png",
    "hc_before_after": r"D:\Something\research\MATE-Automation-V4\10_Papers-Thesis\Physics_Informed\figures\hc_before_after.png",
    "cross_attention_schema": r"D:\Something\research\MATE-Automation-V4\10_Papers-Thesis\Physics_Informed\figures\cross_attention_schema.png",
    "v6_confusion_grid": r"D:\Something\research\MATE-Automation-V4\10_Papers-Thesis\Physics_Informed\figures\v6_confusion_matrix_grid.png",
    "paradigm_comparison": r"D:\Something\research\MATE-Automation-V4\runs\EXP1-XA-HC-100k\20260317_191552_seed42\figures\paradigm_comparison.png",
    "cross_task_comparison": r"D:\Something\research\MATE-Automation-V4\runs\TRK-comparison\figures\cross_task_honest_comparison.png",
    "exp2_accuracy_bars": r"D:\Something\research\MATE-Automation-V4\runs\EXP2-fusion-comparison\figures\exp2_accuracy_bars.png",
    "oedo_sharaq": r"D:\obsidian\OrbitOS\50_Resources\Attachments\OEDO-SHARAQ.png",
    "nst_mate_schematic": r"D:\obsidian\OrbitOS\50_Resources\Physics\literature\MLforPhysics\NST-ML-12C12C-MATE\_page_3_Figure_9.jpeg",
    "nst_tracks": r"D:\obsidian\OrbitOS\50_Resources\Physics\literature\MLforPhysics\NST-ML-12C12C-MATE\_page_5_Figure_3.jpeg",
    "nst_confusion": r"D:\obsidian\OrbitOS\50_Resources\Physics\literature\MLforPhysics\NST-ML-12C12C-MATE\_page_8_Figure_3.jpeg",
    "attention_overlay": r"D:\Something\research\MATE-Automation-V4\10_Papers-Thesis\Physics_Informed\figures\attention_overlay.png",
    "data_pipeline": r"D:\Something\research\MATE-Automation-V4\10_Papers-Thesis\Physics_Informed\figures\data_pipeline.png",
}


def fig_exists(key):
    return os.path.isfile(FIGURES.get(key, ""))


def fig_path(key):
    return FIGURES.get(key, "")


# ── Helpers ────────────────────────────────────────────────────────────────

def add_header_stripe(slide):
    """Add a navy stripe across the top of the slide."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        SLIDE_WIDTH, Inches(1.1)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_BG_STRIPE
    shape.line.fill.background()


def add_slide_number(slide, num, total=5):
    """Add slide number at bottom right."""
    txBox = slide.shapes.add_textbox(
        SLIDE_WIDTH - Inches(1.2), SLIDE_HEIGHT - Inches(0.5),
        Inches(1.0), Inches(0.4)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num} / {total}"
    p.font.size = Pt(11)
    p.font.color.rgb = COLOR_SUBTITLE
    p.alignment = PP_ALIGN.RIGHT


def add_footer_line(slide):
    """Thin accent line near bottom."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), SLIDE_HEIGHT - Inches(0.65),
        SLIDE_WIDTH - Inches(1.0), Pt(2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_ACCENT
    shape.line.fill.background()


def add_title(slide, title_text):
    """Add white title text inside the navy stripe."""
    txBox = slide.shapes.add_textbox(
        Inches(0.7), Inches(0.15),
        SLIDE_WIDTH - Inches(1.4), Inches(0.85)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.font.name = FONT_TITLE
    p.alignment = PP_ALIGN.LEFT


def add_bullets(slide, bullets, left=Inches(0.7), top=Inches(1.35),
                width=None, height=None, font_size=Pt(16), spacing=Pt(6)):
    """Add bullet list. Supports (text, level) tuples or plain strings."""
    if width is None:
        width = SLIDE_WIDTH - Inches(1.4)
    if height is None:
        height = SLIDE_HEIGHT - top - Inches(1.0)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(bullets):
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = text
        p.font.size = font_size
        p.font.color.rgb = COLOR_BODY
        p.font.name = FONT_BODY
        p.level = level
        p.space_after = spacing

        # Highlight key numbers
        if any(kw in text for kw in ['95.80%', '16.1%', '95%', '90-97%',
                                       'p<10', '93%', '97%', '84.6%',
                                       '86.1%', '72.9%']):
            p.font.bold = True

    return txBox


def add_image_safe(slide, key, left, top, max_width=None, max_height=None):
    """Try to embed an image with bounding box constraint; return True if successful."""
    if not fig_exists(key):
        return False
    try:
        from PIL import Image as PILImage
        img = PILImage.open(fig_path(key))
        img_w, img_h = img.size
        aspect = img_w / img_h

        # Start with max dimensions
        w = max_width if max_width else Inches(4.5)
        h = max_height if max_height else Inches(4.8)

        # Fit within bounding box preserving aspect ratio
        if max_width and max_height:
            box_aspect = max_width / max_height
            if aspect > box_aspect:
                # Image is wider than box -> constrain by width
                h = int(max_width / aspect)
            else:
                # Image is taller than box -> constrain by height
                w = int(max_height * aspect)
        elif max_width:
            h = int(max_width / aspect)
        elif max_height:
            w = int(max_height * aspect)

        slide.shapes.add_picture(fig_path(key), left, top, w, h)
        return True
    except Exception as e:
        print(f"  [WARN] Could not embed {key}: {e}")
        return False


def set_notes(slide, notes_text):
    """Set speaker notes."""
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes_text


# ── Build Presentation ────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

# Use blank layout
blank_layout = prs.slide_layouts[6]  # blank

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1: Physics Motivation
# ════════════════════════════════════════════════════════════════════════════
slide1 = prs.slides.add_slide(blank_layout)
add_header_stripe(slide1)
add_title(slide1, "Physics Motivation: Nuclear Structure Through Direct Reactions")
add_footer_line(slide1)
add_slide_number(slide1, 1)

bullets1 = [
    "Nuclear shell structure evolves far from stability",
    ("Traditional magic numbers (2,8,20,28...) change: tensor force reshapes single-particle levels", 1),
    "Inverse-kinematics transfer/knockout reactions probe shell evolution",
    ("Spectroscopic factors C\u00b2S from (d,p), (p,d) at 10\u201330 MeV/u \u2192 J\u03c0, occupancies", 1),
    "Active-target TPCs: ideal for low-energy exotic beams",
    ("Gas = target \u2192 thick-target luminosity without degrading resolution", 1),
    ("Full kinematic coverage, near-unity geometric efficiency for short-range recoils", 1),
    "The bottleneck: reliable particle identification (PID)",
    ("Overlapping dE/dx, pile-up, dynamic range \u2192 extracting physics requires ML-grade PID", 1),
]

# Place bullets on left ~60% of slide
add_bullets(slide1, bullets1, width=Inches(7.5))

# Try to embed MATE schematic on the right
embedded_s1 = add_image_safe(slide1, "nst_mate_schematic",
                              left=Inches(8.3), top=Inches(1.4),
                              max_width=Inches(4.5), max_height=Inches(4.8))

set_notes(slide1, """TALKING POINTS:
- Shell evolution is the fundamental question: how does nuclear structure change as we add neutrons far from stability?
- Magic numbers we learned (2,8,20,28,50,82,126) aren't universal -- tensor force monopole reshuffles levels (Otsuka et al., RMP 2020).
- Direct reactions are the tool: single-step, ~10^-22 s, probe nuclear surface. DWBA analysis extracts C2S.
- Active-target TPCs solve the luminosity problem for exotic beam experiments at facilities like HIRFL-RIBLL (IMP) and RIKEN RIBF.
- But the killer problem is PID: you need clean channel identification to gate on the correct reaction. Classical dE-E cuts struggle with overlapping Bethe-Bloch curves, especially for similar A/Z species.
- This motivates our ML approach.

FIGURE: MATE-TPC schematic from NST paper (pad plane + GEM + field cage).
Figure path: """ + fig_path("nst_mate_schematic"))


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2: The PID Problem in Active-Target TPCs
# ════════════════════════════════════════════════════════════════════════════
slide2 = prs.slides.add_slide(blank_layout)
add_header_stripe(slide2)
add_title(slide2, "The PID Problem in Active-Target TPCs")
add_footer_line(slide2)
add_slide_number(slide2, 2)

bullets2 = [
    "Traditional methods struggle with TPC data",
    ("dE-E: overlapping distributions at similar A/Z (e.g., \u00b3He/\u2074He, \u00b9\u00b3C/\u00b9\u2074C)", 1),
    ("Bragg-curve fitting: pile-up, dynamic range across beam vs. recoil region", 1),
    ("Hough transform mislabels ~1.5% of events (NST-2025-0958)", 1),
    "MATE-TPC at IMP (HIRFL-RIBLL)",
    ("GEM readout, 3,792 triangular pads (48\u00d780), GET electronics", 1),
    ("\u00b9\u00b2C+\u00b9\u00b2C commissioning data: elastic vs. fusion + 5 reaction channels", 1),
    "Key insight: 2D track images encode species-specific spatial signatures",
    ("Bragg peak position, track curvature, energy deposition pattern \u2192 underutilized by cuts", 1),
    "Our approach: CNNs on track images \u2192 the MATE-Automation pipeline",
    ("ResNet-18 backbone + CrossAttention fusion of physics features", 1),
]

add_bullets(slide2, bullets2, width=Inches(7.5))

# Try to embed example tracks on the right -- use nst_tracks (3D views, better aspect ratio)
embedded_s2_tracks = add_image_safe(slide2, "nst_tracks",
                                     left=Inches(8.3), top=Inches(1.35),
                                     max_width=Inches(4.5), max_height=Inches(4.8))
if not embedded_s2_tracks:
    # Fallback to example_tracks with explicit height
    embedded_s2_tracks = add_image_safe(slide2, "example_tracks",
                                         left=Inches(8.3), top=Inches(1.35),
                                         max_width=Inches(4.5), max_height=Inches(4.8))

set_notes(slide2, """TALKING POINTS:
- Traditional PID in TPCs: dE/dx along track + integrated energy. Works well for well-separated species, but struggles for 3He/4He (Delta A/Z ~33%) and especially 13C/14C (Delta A/Z ~8%).
- Hough transform for track finding: our NST paper showed ML corrected ~1.5% mislabeled events.
- MATE detector specifics: GEM amplification (not Micromegas), triangular pad readout (48 rows x 80 pads = 3,792 pads), dual-gain for beam vs recoil.
- The 2D projection images (charge deposition + drift time) contain rich spatial information: Bragg peak location, track length, curvature -- all species-dependent features that a CNN can learn.
- This is why we convert ROOT simulation data -> HDF5 -> 2-channel 80x48 images and feed to ResNet.
- CrossAttention mechanism fuses 4 handcrafted physics features (moments of inertia Iyy, Izz, Iyz, total mass) with the CNN spatial representation.

FIGURE: Example TPC track images showing different particle species.
Figure path (embedded): """ + (fig_path("example_tracks") if embedded_s2_tracks else "NEEDS MANUAL INSERTION") + """
Also useful: """ + fig_path("nst_tracks") + """ (3D track displays from NST paper)
""")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3: Results -- From Event Classification to Trajectory Reconstruction
# ════════════════════════════════════════════════════════════════════════════
slide3 = prs.slides.add_slide(blank_layout)
add_header_stripe(slide3)
add_title(slide3, "Results: Event Classification to Trajectory Reconstruction")
add_footer_line(slide3)
add_slide_number(slide3, 3)

bullets3_left = [
    "PID Classification (NST-2025-0958, \u00b9\u00b2C+\u00b9\u00b2C)",
    ("Elastic/fusion: 90\u201397% accuracy (ResNet-50, sim+exp)", 1),
    ("5-channel fusion discrimination: ~95% (p, \u03b1, n, \u2078Be, 2\u03b1)", 1),
    ("Precision 0.948, Recall 0.914, F1 0.931 on experimental data", 1),
    "Architecture Study (EXP1/EXP2, thesis)",
    ("CrossAtt+HC on 3He/4He: 95.80% (controlled, matched 100k)", 1),
    ("EXP2 fusion comparison: CrossAtt vs Gated vs Concat on 13C/14C", 1),
    "Trajectory Reconstruction (TRK series, new)",
    ("Same framework extended to angle regression", 1),
    ("CrossAtt reduces MAE by 16.1% (0.99\u00b0\u21920.83\u00b0, p<10\u207b\u00b9\u2076)", 1),
    ("High-precision fraction: 65.3%\u219272.9% (+7.6 pp)", 1),
    "Takeaway: pipeline generalizes \u2014 classification AND regression",
]

add_bullets(slide3, bullets3_left, width=Inches(7.2), font_size=Pt(15))

# Try to embed cross-task comparison on the right
embedded_s3 = add_image_safe(slide3, "cross_task_comparison",
                              left=Inches(7.8), top=Inches(1.35),
                              max_width=Inches(5.0), max_height=Inches(4.8))
if not embedded_s3:
    embedded_s3 = add_image_safe(slide3, "paradigm_comparison",
                                  left=Inches(7.8), top=Inches(1.35),
                                  max_width=Inches(5.0), max_height=Inches(4.8))

set_notes(slide3, """TALKING POINTS:
- Three levels of results to present:
  1. Published PID results (NST-2025-0958, submitted Nov 2025): 90-97% elastic/fusion on 12C+12C, ~95% across 5 fusion channels
  2. Thesis architecture study: CrossAtt+HC fusion mechanism. EXP1 showed 95.80% at matched 100k training size. EXP2 compared CrossAtt vs GatedFusion vs ConcatFusion.
  3. NEW: Trajectory reconstruction (TRK series, completed April 2026). Same CrossAttention framework applied to angle regression task.

KEY NUMBERS:
- EXP1: CrossAtt+HC = 95.80% (3He/4He, 100k matched). Architecture effect is +0.024 pp vs baseline (not significant at p=0.813), BUT data scaling to 400k gives +1.77 pp (p<1e-73). This is honest -- architecture helps, but data volume matters more for PID.
- TRK: CrossAtt reduces angle MAE by 16.1% (0.992 -> 0.832 degrees, p<1e-16, CIs fully disjoint). This is the clearest architecture win.
- The generalization story is powerful: same backbone + CrossAtt fusion works for classification AND regression tasks on TPC data.

- 13C/14C is the hardest case: only 8% A/Z difference. CrossAtt gets 84.6% vs ConcatFusion 86.1% -- still challenging.

FIGURES:
- cross_task_honest_comparison.png: Three-panel comparison (Isotope ID / Track Classification / Angle Regression)
  Path: """ + fig_path("cross_task_comparison") + """
- paradigm_comparison.png: Two-paradigm comparison (Range-Energy PID vs ML)
  Path: """ + fig_path("paradigm_comparison") + """
- v6_confusion_matrix_grid.png: Confusion matrices across ablation variants
  Path: """ + fig_path("v6_confusion_grid") + """
""")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4: Connection -- From MATE to DONUTS
# ════════════════════════════════════════════════════════════════════════════
slide4 = prs.slides.add_slide(blank_layout)
add_header_stripe(slide4)
add_title(slide4, "Connection: From MATE-TPC (IMP) to DONUTS (UTokyo CNS)")
add_footer_line(slide4)
add_slide_number(slide4, 4)

bullets4 = [
    "ML-PID for active-target TPCs maps directly to the DONUTS program",
    "OEDO-SHARAQ: inverse-kinematics direct reactions on exotic nuclei",
    ("(d,p), (p,d) transfer reactions at 10\u201330 MeV/u via RIKEN RIBF", 1),
    ("Same PID challenge: separate light recoils (p, d, \u03b1) from beam background", 1),
    ("DG-M-THGEM active target: GEM-based, dual gain, stable to 2.5\u00d710\u2076 pps", 1),
    "SAKURA: surrogate (d,p) for r-process (n,\u03b3) cross-sections",
    ("\u2079\u2079Se(d,p) flagship result (Imai et al., PLB 2024) \u2192 reliable PID essential", 1),
    ("Extending to more exotic targets: ML-PID improves channel discrimination", 1),
    "Transferable capabilities from my work:",
    ("ResNet + CrossAtt framework adaptable to different TPC geometries", 1),
    ("Full pipeline: simulation \u2192 image \u2192 ML inference \u2192 validated on experimental data", 1),
    ("Track reconstruction + angle regression adds trajectory analysis capability", 1),
]

add_bullets(slide4, bullets4, width=Inches(7.5))

# Try to embed OEDO-SHARAQ image on the right
embedded_s4 = add_image_safe(slide4, "oedo_sharaq",
                              left=Inches(8.3), top=Inches(1.5),
                              max_width=Inches(4.5), max_height=Inches(4.5))

set_notes(slide4, """TALKING POINTS:
- This slide is the BRIDGE -- connecting IMP research to Imai's program at UTokyo CNS.
- DONUTS group operates at RIKEN RIBF via OEDO-SHARAQ beamline: energy-degraded exotic beams for direct reactions.
- The PID challenge is identical in principle: separate light recoils from heavy beam background in an active-target TPC environment.
- SAKURA program: surrogate (d,p) reactions to constrain (n,gamma) for r-process nucleosynthesis. The 79Se(d,p) result (Imai, PLB 2024) is the flagship -- shows why clean PID matters for cross-section extraction.
- DG-M-THGEM detector (Iwamoto, Ota, Imai et al., PTEP 2023): dual-gain GEM-based active target, similar in principle to MATE. ML-PID framework transfers naturally.
- Key message: I'm not asking to replicate my thesis -- I'm bringing a validated toolkit that addresses a real analysis bottleneck in DONUTS experiments.
- Mention CAT-M detector (shared CNS infrastructure between NUSPEQ/Aoi and DONUTS/Imai groups).

FIGURE: OEDO-SHARAQ beamline schematic.
Path: """ + fig_path("oedo_sharaq") + """
""")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5: Outlook
# ════════════════════════════════════════════════════════════════════════════
slide5 = prs.slides.add_slide(blank_layout)
add_header_stripe(slide5)
add_title(slide5, "Outlook: Next Steps and Open Questions")
add_footer_line(slide5)
add_slide_number(slide5, 5)

bullets5 = [
    "Extend CrossAtt-PID framework to DONUTS detector data",
    ("Adapt pipeline from MATE triangular pads \u2192 DG-M-THGEM geometry", 1),
    ("Leverage simulation infrastructure for new target/beam combinations", 1),
    "Open questions for discussion",
    ("Transfer learning: how well does the model generalize across TPC geometries?", 1),
    ("Beam conditions: different rates, gas mixtures, dynamic range at RIKEN RIBF", 1),
    ("Complex reaction channels in SAKURA: multi-body final states, (d,p\u03b3) coincidences", 1),
    ("Integration with existing ATPC Flow analysis framework (Python/C++/Rust)", 1),
    "Broader vision",
    ("ML-augmented analysis as standard tool for next-generation RI-beam experiments", 1),
    ("From PID \u2192 full event reconstruction: vertex, track, angle, channel ID in one pipeline", 1),
]

add_bullets(slide5, bullets5, width=Inches(11.0))

# Add a subtle "Thank you / Discussion" at the bottom center
txBox = slide5.shapes.add_textbox(
    Inches(3.5), SLIDE_HEIGHT - Inches(1.4),
    Inches(6.0), Inches(0.6)
)
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Looking forward to your guidance on the highest-impact application."
p.font.size = Pt(16)
p.font.italic = True
p.font.color.rgb = COLOR_ACCENT
p.font.name = FONT_BODY
p.alignment = PP_ALIGN.CENTER

set_notes(slide5, """TALKING POINTS:
- This is the forward-looking slide. End on a collaborative note -- invite Imai's input.
- Key question to ASK: "Given the DONUTS program's current priorities, where would ML-based PID have the most immediate impact?"
- Transfer learning is a genuine open question: MATE has 3,792 triangular pads; DONUTS detectors may have different geometries. Domain adaptation (DANN/MCD) methods are ready to explore.
- SAKURA reactions are more complex than simple binary classification -- multi-body final states may need multi-class or hierarchical classification approaches.
- Mention: the MATE-Automation codebase is modular and well-documented. Adapting to a new detector is a config change + retraining, not a rewrite.
- The ATPC Flow framework (mentioned in RCNP collaboration meeting notes) is the existing analysis software. Integration point: ML inference as a module within that pipeline.
- End with: "I'm eager to learn how this work can contribute to the group's goals, and I welcome your suggestions on the most valuable direction to pursue during the PhD."

TIMELINE CONTEXT (if asked):
- IMP graduation: June 2027
- UTokyo PhD D1 enrollment: October 2027 (primary) or October 2028 (fallback via kenkyusei)
- TOEFL iBT Attempt #1: August 30, 2026
- GRE Physics: Sept/Oct 2026
- UTokyo entrance exam: Jan-Feb 2027 (estimated)
""")


# ── Save ───────────────────────────────────────────────────────────────────
os.makedirs(OUT_DIR, exist_ok=True)
prs.save(OUT_FILE)
print(f"Saved: {OUT_FILE}")

# Report which figures were embedded
print("\n-- Figure embedding report --")
for key, path in FIGURES.items():
    status = "EXISTS" if os.path.isfile(path) else "MISSING"
    print(f"  {key}: {status}")
